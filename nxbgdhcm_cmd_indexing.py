import os
import sys
import time
import json
import re
import hashlib
import uuid
import platform
import getpass
import requests
import base64
import threading
import signal
from datetime import datetime
from pdf2image import convert_from_path
from PIL import Image
import fitz

# Bật hỗ trợ màu ANSI cho Windows CMD/PowerShell
os.system('') 

class Colors:
    SYSTEM = '\033[93m'      # Vàng
    SUCCESS = '\033[92m'     # Xanh lá
    ERROR = '\033[91m'       # Đỏ
    INFO = '\033[97m'        # Trắng
    AI_THINK = '\033[94m'    # Xanh dương sáng
    AI_STREAM = '\033[32m'   # Xanh lá cây đậm
    WARNING = '\033[95m'     # Hồng tím
    ANIM = '\033[93m'        # Vàng
    RESET = '\033[0m'

# Thử nạp các thư viện Office và MySQL
try:
    import docx
except ImportError:
    docx = None
try:
    import pandas as pd
except ImportError:
    pd = None
try:
    import pptx
except ImportError:
    pptx = None
try:
    import win32com.client as win32
    import pythoncom
except ImportError:
    win32 = None
    pythoncom = None
    
try:
    import mysql.connector
    HAS_MYSQL = True
except ImportError:
    HAS_MYSQL = False

# ==========================================
# BIẾN TOÀN CỤC & CẤU HÌNH
# ==========================================
API_URL = ""
MODEL = ""
API_KEY = ""
HEADERS = {}

is_running = True 

if getattr(sys, 'frozen', False):
    BASE_DIR = os.path.dirname(sys.executable)
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

POPPLER_PATH = os.path.join(BASE_DIR, "poppler", "Library", "bin")
TEMP_DIR = os.path.join(BASE_DIR, "AI_Temp")
SQL_INI_FILE = os.path.join(BASE_DIR, "sql.ini")

APPDATA_DIR = os.path.join(os.getenv('APPDATA'), "NXBGDHCM_OCR")
if not os.path.exists(APPDATA_DIR):
    os.makedirs(APPDATA_DIR)
KEY_FILE = os.path.join(APPDATA_DIR, "nhanviennxbgdhcm.json")

LOGS_DIR = os.path.join(BASE_DIR, "logs")
if not os.path.exists(LOGS_DIR):
    os.makedirs(LOGS_DIR)
SESSION_LOG_FILE = os.path.join(LOGS_DIR, datetime.now().strftime("%y%m%d-%H%M%S.log"))

DB_USER = ""
DB_PASS = ""
DB_HOST = ""
DB_PORT = 3306
DB_NAME = ""
PERSON_KEY = ""

# ==========================================
# HÀM BẮT SỰ KIỆN CTRL+C (GRACEFUL EXIT)
# ==========================================
def signal_handler(sig, frame):
    global is_running
    if is_running:
        is_running = False
        log("[DUNG] Đã nhận lệnh dừng (Ctrl+C)! Đang hoàn tất file cuối trước khi thoát...", "warning")
    else:
        log("[CANH BAO] Ép buộc thoát ngay lập tức!", "error")
        sys.exit(0)

signal.signal(signal.SIGINT, signal_handler)

# ==========================================
# HÀM LOG CLI VÀ GHI FILE JSON
# ==========================================
def log(message, level="info"):
    now = datetime.now()
    console_time = now.strftime("[%y-%m-%d] [%H:%M:%S]")
    json_time = now.strftime("%y-%m-%d :: %H:%M:%S")
    
    color = Colors.INFO
    if level == "system": color = Colors.SYSTEM
    elif level == "success": color = Colors.SUCCESS
    elif level == "error": color = Colors.ERROR
    elif level == "warning": color = Colors.WARNING
    elif level == "ai": color = Colors.AI_THINK
    
    print(f"{color}{console_time} {message}{Colors.RESET}")
    
    if level != "ai":
        try:
            log_entry = {
                "thời điểm": json_time,
                "Thông báo": message
            }
            with open(SESSION_LOG_FILE, "a", encoding="utf-8") as f:
                f.write(json.dumps(log_entry, ensure_ascii=False) + "\n")
        except Exception:
            pass

# ==========================================
# HÀM XỬ LÝ MYSQL (DÙNG CHUNG)
# ==========================================
def init_person_key():
    global PERSON_KEY
    mac_num = uuid.getnode()
    mac_address = ':'.join(('%012X' % mac_num)[i:i+2] for i in range(0, 12, 2))
    computer_name = platform.node()
    user_name = getpass.getuser()

    if os.path.exists(KEY_FILE):
        try:
            with open(KEY_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                PERSON_KEY = data.get("Person_key", "")
        except: pass

    if not PERSON_KEY:
        PERSON_KEY = str(uuid.uuid4())
        with open(KEY_FILE, "w", encoding="utf-8") as f:
            json.dump({"Person_key": PERSON_KEY, "Mac_Address": mac_address}, f, indent=4)

    return mac_address, computer_name, user_name

def init_sql_config():
    global DB_HOST, DB_PORT, DB_USER, DB_PASS, DB_NAME
    if not os.path.exists(SQL_INI_FILE):
        with open(SQL_INI_FILE, "w", encoding="utf-8") as f:
            f.write("ip: 192.168.192.12\nport: 3306\nuser: nxbgdhcm\npass: chitan1811\ndatabase: nxbgdhcm_vanban\n")
    try:
        with open(SQL_INI_FILE, "r", encoding="utf-8") as f:
            for line in f:
                if ":" in line:
                    k, v = line.split(":", 1)
                    k = k.strip().lower()
                    v = v.strip()
                    if k == "ip": DB_HOST = v
                    elif k == "port": DB_PORT = int(v) if v.isdigit() else 3306
                    elif k == "user": DB_USER = v
                    elif k == "pass": DB_PASS = v
                    elif k == "database": DB_NAME = v
    except Exception as e:
        log(f"[CANH BAO] Lỗi đọc file sql.ini: {e}", "warning")

def get_db_connection():
    if not HAS_MYSQL: raise Exception("Thư viện mysql-connector-python chưa được cài đặt!")
    return mysql.connector.connect(host=DB_HOST, port=DB_PORT, user=DB_USER, password=DB_PASS, database=DB_NAME)

def init_global_db():
    if not os.path.exists(TEMP_DIR): os.makedirs(TEMP_DIR)
    mac_addr, comp_name, os_user = init_person_key()
    init_sql_config()
    
    if not HAS_MYSQL: return
    try:
        conn_init = mysql.connector.connect(host=DB_HOST, port=DB_PORT, user=DB_USER, password=DB_PASS)
        c_init = conn_init.cursor()
        c_init.execute(f"CREATE DATABASE IF NOT EXISTS {DB_NAME} CHARACTER SET utf8mb4 COLLATE utf8mb4_unicode_ci")
        conn_init.commit()
        c_init.close()
        conn_init.close()
        
        conn = get_db_connection()
        c = conn.cursor()
        c.execute('''CREATE TABLE IF NOT EXISTS Nguoi_dung (
                        Person_key VARCHAR(255) PRIMARY KEY, Mac_Address VARCHAR(100),
                        Computer_Name VARCHAR(255), User_Name VARCHAR(255),
                        Ho_Va_Ten VARCHAR(255) DEFAULT '', Chuc_vu VARCHAR(255) DEFAULT '', Phong_Ban VARCHAR(255) DEFAULT ''
                     )''')
        c.execute('''INSERT INTO Nguoi_dung (Person_key, Mac_Address, Computer_Name, User_Name) 
                     VALUES (%s, %s, %s, %s) ON DUPLICATE KEY UPDATE 
                     Mac_Address=VALUES(Mac_Address), Computer_Name=VALUES(Computer_Name), User_Name=VALUES(User_Name)
                  ''', (PERSON_KEY, mac_addr, comp_name, os_user))
        
        c.execute('''CREATE TABLE IF NOT EXISTS documents (
                        md5 VARCHAR(255), person_key VARCHAR(255), file_name TEXT, file_path TEXT,
                        file_type VARCHAR(50), doc_type VARCHAR(255), doc_number VARCHAR(255),
                        doc_day VARCHAR(10), doc_month VARCHAR(10), doc_year VARCHAR(10),
                        doc_org TEXT, doc_signer TEXT, full_text LONGTEXT, last_scan BIGINT,
                        PRIMARY KEY (md5, person_key)
                     )''')
        c.execute('''CREATE TABLE IF NOT EXISTS deleted_documents (
                        md5 VARCHAR(255), person_key VARCHAR(255), file_name TEXT, file_path TEXT,
                        file_type VARCHAR(50), doc_type VARCHAR(255), doc_number VARCHAR(255),
                        doc_day VARCHAR(10), doc_month VARCHAR(10), doc_year VARCHAR(10),
                        doc_org TEXT, doc_signer TEXT, full_text LONGTEXT, deleted_time BIGINT,
                        PRIMARY KEY (md5, person_key)
                     )''')
                     
        c.execute('''CREATE TABLE IF NOT EXISTS ket_noi_ai (
                        ID INT AUTO_INCREMENT PRIMARY KEY, Preset_Name VARCHAR(255) NOT NULL,
                        URL VARCHAR(255) NOT NULL, Model_Name VARCHAR(255) NOT NULL,
                        API_Key VARCHAR(255) NOT NULL, `Default` BOOLEAN DEFAULT FALSE
                     )''')
                     
        c.execute('''CREATE TABLE IF NOT EXISTS loai_van_ban (
                        ID INT AUTO_INCREMENT PRIMARY KEY, Loai_VB VARCHAR(255) NOT NULL UNIQUE
                     )''')

        # Khởi tạo tự động bảng quản lý prompt tập trung
        c.execute('''CREATE TABLE IF NOT EXISTS cau_hinh_prompt (
                        Prompt_Key VARCHAR(50) PRIMARY KEY, Prompt_Content TEXT NOT NULL,
                        prompt_type VARCHAR(50) NOT NULL DEFAULT 'in_app', Description VARCHAR(255) DEFAULT NULL
                     ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4 COLLATE=utf8mb4_unicode_ci;''')

        conn.commit()
        c.close()
        conn.close()
    except Exception as e:
        log(f"[LOI MYSQL] Khởi tạo thất bại: {e}", "error")
        sys.exit(1)

def calculate_md5(filepath):
    hash_md5 = hashlib.md5()
    try:
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""): hash_md5.update(chunk)
        return hash_md5.hexdigest()
    except Exception: return None

def check_and_update_db(db_conn, filepath, filename, current_md5):
    try:
        db_conn.ping(reconnect=True, attempts=3, delay=1)
        c = db_conn.cursor()
        current_time = int(time.time())
        
        c.execute("SELECT md5 FROM documents WHERE md5 = %s AND person_key = %s", (current_md5, PERSON_KEY))
        if c.fetchone():
            c.execute('''UPDATE documents SET file_name = %s, file_path = %s, last_scan = %s 
                         WHERE md5 = %s AND person_key = %s''', (filename, filepath, current_time, current_md5, PERSON_KEY))
            db_conn.commit(); c.close()
            return True
            
        c.execute('''SELECT md5, person_key, file_name, file_path, file_type, 
                            doc_type, doc_number, doc_day, doc_month, doc_year, 
                            doc_org, doc_signer, full_text 
                     FROM deleted_documents WHERE md5 = %s AND person_key = %s''', (current_md5, PERSON_KEY))
        row = c.fetchone()
        if row:
            c.execute('''INSERT INTO documents 
                         (md5, person_key, file_name, file_path, file_type, doc_type, doc_number, doc_day, doc_month, doc_year, doc_org, doc_signer, full_text, last_scan)
                         VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)''',
                      (*row, current_time))
            c.execute("DELETE FROM deleted_documents WHERE md5 = %s AND person_key = %s", (current_md5, PERSON_KEY))
            db_conn.commit(); c.close()
            return True
        c.close()
        return False
    except Exception: return False

def save_document_to_db(db_conn, md5_hash, filename, filepath, file_type, metadata, full_text):
    try:
        db_conn.ping(reconnect=True, attempts=3, delay=1)
        c = db_conn.cursor()
        c.execute('''INSERT INTO documents 
                     (md5, person_key, file_name, file_path, file_type, doc_type, doc_number, doc_day, doc_month, doc_year, doc_org, doc_signer, full_text, last_scan)
                     VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                     ON DUPLICATE KEY UPDATE
                     file_name=VALUES(file_name), file_path=VALUES(file_path), file_type=VALUES(file_type),
                     doc_type=VALUES(doc_type), doc_number=VALUES(doc_number), doc_day=VALUES(doc_day), 
                     doc_month=VALUES(doc_month), doc_year=VALUES(doc_year), doc_org=VALUES(doc_org), 
                     doc_signer=VALUES(doc_signer), full_text=VALUES(full_text), last_scan=VALUES(last_scan)
                     ''',
                  (md5_hash, PERSON_KEY, filename, filepath, file_type,
                   metadata.get("Loại văn bản", ""), metadata.get("Số văn bản", ""), metadata.get("Ngày", ""), metadata.get("Tháng", ""),
                   metadata.get("Năm", ""), metadata.get("Đơn vị soạn văn bản", ""), metadata.get("Người ký", ""), full_text, int(time.time())))
        db_conn.commit()
        c.close()
    except Exception as e:
        log(f"[LOI GHI CSDL]: {e}", "error")

# ==========================================
# CẤU HÌNH AI VÀ PROMPT DYNAMIC FETCH
# ==========================================
def load_connections():
    global API_URL, MODEL, API_KEY, HEADERS
    if not HAS_MYSQL:
        log("[LOI] Hệ thống cần kết nối MySQL để tải cấu hình AI!", "error")
        sys.exit(1)
        
    try:
        conn = get_db_connection()
        c = conn.cursor()
        c.execute("SELECT Preset_Name, URL, Model_Name, API_Key, `Default` FROM ket_noi_ai")
        rows = c.fetchall()
        c.close()
        conn.close()

        if not rows:
            log("[LOI] Bảng ket_noi_ai trên Database đang trống! Vui lòng cấu hình AI trước.", "error")
            sys.exit(1)

        selected = next((r for r in rows if r[4]), rows[0])
        name, API_URL, MODEL, API_KEY, _ = selected
        HEADERS = {"Content-Type": "application/json", "Authorization": f"Bearer {API_KEY}"}
        log(f"[KET NOI] Đã thiết lập liên kết Server: [{name}]", "success")
        
    except Exception as e: 
        log(f"[LOI] Tải cấu hình AI từ Database thất bại: {e}", "error")
        sys.exit(1)

def load_document_types(db_conn=None):
    types = []
    if HAS_MYSQL and db_conn:
        try:
            c = db_conn.cursor()
            c.execute("SELECT Loai_VB FROM loai_van_ban")
            rows = c.fetchall()
            types = [row[0].strip() for row in rows if row[0]]
            c.close()
            if types:
                return ", ".join(types)
        except Exception as e:
            log(f"[CANH BAO] Lỗi tải danh mục loại văn bản từ DB: {e}", "warning")
    return ""

def get_prompt_from_db(db_conn, prompt_key):
    try:
        db_conn.ping(reconnect=True, attempts=3, delay=1)
        c = db_conn.cursor()
        c.execute("SELECT Prompt_Content FROM cau_hinh_prompt WHERE Prompt_Key = %s", (prompt_key,))
        row = c.fetchone()
        c.close()
        return row[0] if row else ""
    except:
        return ""

def get_prompt_page_1(db_conn, doc_types_str=""):
    prompt = get_prompt_from_db(db_conn, 'prompt_ocr_page_1')
    if not prompt:
        prompt = 'Bạn là một hệ thống trích xuất dữ liệu API tự động. BẠN PHẢI TRẢ VỀ ĐÚNG MỘT KHỐI JSON DUY NHẤT. TUYỆT ĐỐI KHÔNG giải thích.\n- Hãy trích xuất các nội dung sau theo cấu trúc Json: "Loại văn bản" ({types_instruction}nếu không rõ ghi "Văn bản"); "Số văn bản"; "Ngày" (2 số); "Tháng" (2 số); "Năm" (4 số); "Đơn vị soạn văn bản"; "Người ký" (Nếu nhiều người thì cách nhau bằng dấu phẩy); "Toàn văn" (nếu là ảnh scan của văn bản, trả về nội dung toàn văn, nếu là ảnh chụp/bản vẽ, trả về nội dung phân tích chi tiết bức ảnh).\nLƯU Ý QUAN TRỌNG: 1. Kết quả bắt buộc phải bắt đầu bằng { và kết thúc bằng }. 2. Luôn đặt trường "Toàn văn" ở cuối cùng. 3. Nếu không tìm thấy thông tin nào, để chuỗi rỗng "".'
    types_instruction = f'Dựa vào danh sách sau: [{doc_types_str}]. ' if doc_types_str else ''
    return prompt.replace('{types_instruction}', types_instruction)

def get_prompt_word_metadata(db_conn, doc_types_str=""):
    prompt = get_prompt_from_db(db_conn, 'prompt_text_metadata')
    if not prompt:
        prompt = 'Bạn là một hệ thống trích xuất dữ liệu API tự động. BẠN PHẢI TRẢ VỀ ĐÚNG MỘT KHỐI JSON DUY NHẤT.\nHãy trích xuất các thông tin sau theo cấu trúc Json: "Loại văn bản" ({types_instruction}nếu không rõ ghi "Văn bản"); "Số văn bản"; "Ngày" (2 số); "Tháng" (2 số); "Năm" (4 số); "Đơn vị soạn văn bản"; "Người ký".\nLƯU Ý QUAN TRỌNG: 1. Bắt buộc bắt đầu bằng { và kết thúc bằng }. 2. KHÔNG trích xuất "Toàn văn".\n\n[NỘI DUNG VĂN BẢN]:\n'
    types_instruction = f'Dựa vào danh sách sau: [{doc_types_str}]. ' if doc_types_str else ''
    return prompt.replace('{types_instruction}', types_instruction)

PROMPT_PAGE_N = """Bạn là một hệ thống trích xuất dữ liệu API tự động. BẠN PHẢI TRẢ VỀ ĐÚNG MỘT KHỐI JSON DUY NHẤT.
Trích xuất nội dung của trang này theo cấu trúc Json chỉ với 2 trường: "Người ký"; "Toàn văn" (nếu là ảnh scan của văn bản, trả về nội dung toàn văn, nếu là ảnh chụp/bản vẽ, trả về nội dung phân tích chi tiết bức ảnh).
LƯU Ý QUAN TRỌNG: 
1. Luôn đặt trường "Toàn văn" ở cuối cùng trong khối JSON.
2. Nếu gặp chữ ký không đọc được, ghi "[Ký tên và Đóng dấu]". """

def extract_json(text, is_page_1=False):
    try: return json.loads(text, strict=False)
    except: pass
    match = re.search(r'`{3}(?:json)?(.*?)`{3}', text, re.DOTALL | re.IGNORECASE)
    if match:
        try: return json.loads(match.group(1).strip(), strict=False)
        except: pass
    start, end = text.find('{'), text.rfind('}')
    if start != -1 and end != -1 and end > start:
        try: return json.loads(text[start:end+1], strict=False)
        except: pass
        
    result = {}
    signer_match = re.search(r'"Người ký"\s*:\s*"?(.*?)"?(?:,|"Toàn văn"|\Z)', text, re.IGNORECASE | re.DOTALL)
    if signer_match: result["Người ký"] = signer_match.group(1).strip().strip('"').strip()
    tv_match = re.search(r'"Toàn văn"\s*:\s*(.*)', text, re.IGNORECASE | re.DOTALL)
    if tv_match:
        tv_text = tv_match.group(1).strip()
        if tv_text.startswith('"'): tv_text = tv_text[1:]
        if tv_text.endswith('}'): tv_text = tv_text[:-1].strip()
        if tv_text.endswith('"'): tv_text = tv_text[:-1].strip()
        result["Toàn văn"] = tv_text
    if is_page_1:
        for key in ["Loại văn bản", "Số văn bản", "Ngày", "Tháng", "Năm", "Đơn vị soạn văn bản"]:
            m = re.search(f'"{key}"\\s*:\\s*"?(.*?)"?(?:,|$)', text, re.IGNORECASE)
            if m: result[key] = m.group(1).strip().strip('"').strip()
    
    if result.get("Toàn văn") or result.get("Người ký") or (is_page_1 and result.get("Loại văn bản")):
        return result
    return None

def merge_signers(existing_signers, new_signers):
    if not new_signers: return existing_signers
    if not existing_signers: return new_signers
    existing_list = [s.strip() for s in re.split(r'[,;]', existing_signers) if s.strip()]
    new_list = [s.strip() for s in re.split(r'[,;]', new_signers) if s.strip()]
    for s in new_list:
        if s not in existing_list: existing_list.append(s)
    return ", ".join(existing_list)

# ==========================================
# CƠ CHẾ GỌI API & HOẠT ẢNH TERMINAL
# ==========================================
def call_ai_stream(prompt_text, base64_img, page_num, total_pages, filename, file_prog_str, force_text_mode=False):
    global is_running
    messages = []
    if base64_img and not force_text_mode:
        messages.append({"role": "user", "content": [{"type": "text", "text": prompt_text}, {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_img}"}}]})
    else:
        messages.append({"role": "user", "content": prompt_text})

    payload = {"model": MODEL, "messages": messages, "stream": True, "temperature": 0.1, "max_tokens": 4000}
    full_response = ""
    
    now_str = datetime.now().strftime("[%y-%m-%d] [%H:%M:%S]")
    if force_text_mode:
        msg = f"{Colors.AI_THINK}[A.I] đang xử lý khối {Colors.SUCCESS}{page_num}/{total_pages}{Colors.AI_THINK} của file {Colors.SYSTEM}{filename}{Colors.AI_THINK} ({Colors.WARNING}{file_prog_str}{Colors.AI_THINK}) "
    else:
        msg = f"{Colors.AI_THINK}[A.I] đang xử lý trang {Colors.SUCCESS}{page_num}/{total_pages}{Colors.AI_THINK} của file {Colors.SYSTEM}{filename}{Colors.AI_THINK} ({Colors.WARNING}{file_prog_str}{Colors.AI_THINK}) "
        
    print(f"\n{Colors.INFO}{now_str}{Colors.RESET} {msg}", end="", flush=True)

    is_waiting = True
    def light_anim():
        frames = ["▓▒░░░", "▒▓▒░░", "░▒▓▒░", "░░▒▓▒", "░░░▒▓", "░░▒▓▒", "░▒▓▒░", "▒▓▒░░"]
        i = 0
        while is_waiting and is_running:
            print(f"{Colors.ANIM}{frames[i % len(frames)]}{Colors.RESET}", end="", flush=True)
            time.sleep(0.15)
            if is_waiting and is_running:
                print("\b" * 5, end="", flush=True)
            i += 1
        print("\b" * 5 + " " * 5 + "\b" * 5, end="", flush=True)
    
    anim_thread = threading.Thread(target=light_anim)
    anim_thread.start()
    
    try:
        with requests.post(API_URL, headers=HEADERS, json=payload, stream=True, timeout=300) as response:
            response.raise_for_status()
            first_chunk = False
            
            for line in response.iter_lines():
                if not is_running: break
                if line:
                    decoded_line = line.decode('utf-8').strip()
                    if decoded_line.startswith("data: "):
                        data_str = decoded_line[6:]
                        if data_str == "[DONE]": break
                        try:
                            data = json.loads(data_str)
                            chunk = data.get("choices", [{}])[0].get("delta", {}).get("content", "")
                            if chunk:
                                if not first_chunk:
                                    is_waiting = False
                                    anim_thread.join(timeout=2)
                                    if is_running: 
                                        print(f"\n{Colors.AI_STREAM}   => ", end="", flush=True)
                                    first_chunk = True
                                full_response += chunk
                                if is_running: 
                                    print(f"{Colors.AI_STREAM}{chunk}", end="", flush=True)
                        except json.JSONDecodeError: pass
    except Exception as e:
        is_waiting = False
        print(f"\n{Colors.INFO}{now_str}{Colors.RESET} {Colors.ERROR}[LOI] Lỗi kết nối AI: {e}{Colors.RESET}")
        return ""
        
    is_waiting = False
    print(Colors.RESET) 
    return full_response

# ==========================================
# TIẾN TRÌNH LÕI (CLI BATCH PROCESS)
# ==========================================
def get_all_target_files(folder_path):
    target_files = []
    valid_exts = ['.pdf', '.jpg', '.jpeg', '.png', '.doc', '.docx', '.xls', '.xlsx', '.txt', '.ppt', '.pptx']
    
    for root_dir, dirs, files in os.walk(folder_path):
        for file in files:
            if file.startswith('~$') or file.startswith('._'):
                continue
                
            ext = os.path.splitext(file)[1].lower()
            if ext in valid_exts:
                filepath = os.path.join(root_dir, file)
                try:
                    if os.path.getsize(filepath) == 0:
                        continue
                    target_files.append(filepath)
                except Exception:
                    pass
                    
    return target_files

def process_batch(folder_path):
    global is_running
    if pythoncom: pythoncom.CoInitialize()
    
    db_conn = None
    if HAS_MYSQL:
        try:
            db_conn = get_db_connection()
        except Exception as e:
            log(f"[LOI] Lỗi khởi tạo Database tổng: {e}", "error")
            return

    try:
        all_files = get_all_target_files(folder_path)
        total_files = len(all_files)
        if total_files == 0:
            log("[CANH BAO] Không tìm thấy tập tin nào hợp lệ trong thư mục này.", "warning")
            return

        log(f"[QUET] Bắt đầu phân tích {total_files} tập tin. Đang tính toán MD5...", "info")
        file_data = []
        for i, filepath in enumerate(all_files):
            if not is_running: break
            filename = os.path.basename(filepath)
            md5_val = calculate_md5(filepath)
            if md5_val: file_data.append({"path": filepath, "name": filename, "md5": md5_val})
            
            now_str = datetime.now().strftime("[%y-%m-%d] [%H:%M:%S]")
            print(f"\r{Colors.INFO}{now_str}{Colors.RESET} {Colors.WARNING}[QUET MD5] Tiến độ: {i+1}/{total_files}{Colors.RESET}", end="")
        print() 
        
        if not is_running: return
        log("[THANH CONG] Quét MD5 hoàn tất. Bắt đầu đối chiếu và trích xuất!", "success")
        
        doc_types_str = load_document_types(db_conn)
        
        success_count = 0
        skip_count = 0
        failed_files = []

        def run_files(files_list, is_retry=False):
            nonlocal success_count, skip_count
            
            # Đọc động prompt cho trang n tự động từ CSDL
            db_prompt_page_n = get_prompt_from_db(db_conn, 'prompt_cmd_ocr_page_n')
            current_prompt_page_n = db_prompt_page_n if db_prompt_page_n else PROMPT_PAGE_N
            
            for data in files_list:
                if not is_running: break
                filepath, filename, file_md5 = data["path"], data["name"], data["md5"]
                ext = os.path.splitext(filename)[1].lower()
                is_excel = ext in ['.xls', '.xlsx']
                is_word = ext in ['.doc', '.docx']
                is_ppt = ext in ['.ppt', '.pptx']
                is_txt = ext == '.txt'
                file_type_str = ext.replace(".", "").upper()
                
                current_processed = success_count + skip_count
                file_prog_str = f"{current_processed+1}/{total_files}"
                
                log(f"\n{'='*60}", "system")
                log(f"[FILE] Đang xử lý: {filename} [{file_prog_str}]", "system")

                if not is_retry and check_and_update_db(db_conn, filepath, filename, file_md5):
                    log(f"[SKIP] Đã có trong Database. Bỏ qua AI (Skip/Upsert).", "success")
                    skip_count += 1
                    continue

                if is_txt:
                    full_text = ""
                    try:
                        log("[TEXT] Đang đọc nội dung file văn bản...", "info")
                        try:
                            with open(filepath, 'r', encoding='utf-8') as f:
                                full_text = f.read()
                        except UnicodeDecodeError:
                            with open(filepath, 'r', encoding='windows-1258') as f:
                                full_text = f.read()
                    except Exception as e:
                        log(f"[LOI] Lỗi đọc file TXT: {e}", "error")
                        continue 

                    if is_running:
                        meta = {"Loại văn bản": "Văn bản thuần túy (TXT)", "Số văn bản": "", "Ngày": "", "Tháng": "", "Năm": "", "Đơn vị soạn văn bản": "", "Người ký": ""}
                        save_document_to_db(db_conn, file_md5, filename, filepath, file_type_str, meta, full_text.strip())
                        log(f"[THANH CONG] Hoàn thành nạp DB Text.", "success")
                        if not is_retry: success_count += 1
                    continue

                elif is_ppt:
                    full_text = ""
                    try:
                        log("[POWERPOINT] Đang bóc tách Text từ Slide...", "info")
                        if ext == '.pptx':
                            if not pptx:
                                log("[LOI] Chưa cài đặt thư viện python-pptx.", "error")
                                continue 
                            prs = pptx.Presentation(filepath)
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        full_text += shape.text + "\n"
                        elif ext == '.ppt':
                            if not win32:
                                log("[LOI] Chưa cài đặt thư viện pywin32.", "error")
                                continue 
                            ppt_app = None
                            prs = None
                            try:
                                ppt_app = win32.Dispatch("PowerPoint.Application")
                                prs = ppt_app.Presentations.Open(os.path.abspath(filepath), ReadOnly=True, WithWindow=False)
                                for slide in prs.Slides:
                                    for shape in slide.Shapes:
                                        if shape.HasTextFrame:
                                            if shape.TextFrame.HasText:
                                                full_text += shape.TextFrame.TextRange.Text + "\n"
                            finally:
                                if prs:
                                    try: prs.Close()
                                    except: pass
                                if ppt_app:
                                    try: ppt_app.Quit()
                                    except: pass
                    except Exception as e:
                        log(f"[LOI] Lỗi PowerPoint: {e}", "error")
                        continue 

                    if is_running:
                        meta = {"Loại văn bản": "Tài liệu Trình chiếu", "Số văn bản": "", "Ngày": "", "Tháng": "", "Năm": "", "Đơn vị soạn văn bản": "", "Người ký": ""}
                        save_document_to_db(db_conn, file_md5, filename, filepath, file_type_str, meta, full_text.strip())
                        log(f"[THANH CONG] Hoàn thành nạp DB PowerPoint.", "success")
                        if not is_retry: success_count += 1
                    continue

                elif is_excel:
                    if not pd:
                        log("[LOI] Chưa cài đặt thư viện Pandas.", "error")
                        continue 
                    try:
                        log("[EXCEL] Đang convert Excel -> CSV...", "info")
                        xls = pd.ExcelFile(filepath)
                        full_csv_text = ""
                        for sheet in xls.sheet_names:
                            if not is_running: break
                            df = pd.read_excel(xls, sheet_name=sheet)
                            full_csv_text += f"===== sheet: {sheet} =====\n{df.to_csv(index=False)}\n\n"
                        if is_running:
                            save_document_to_db(db_conn, file_md5, filename, filepath, file_type_str, {"Loại văn bản": "Bảng tính Excel"}, full_csv_text)
                            log(f"[THANH CONG] Hoàn tất nạp Excel.", "success")
                            if not is_retry: success_count += 1
                    except Exception as e:
                        log(f"[LOI] Lỗi Excel: {e}", "error")
                    continue

                elif is_word:
                    full_text = ""
                    try:
                        log("[WORD] Đang bóc tách RAW Text và Bảng biểu từ Word...", "info")
                        if ext == '.docx':
                            if not docx: raise Exception("Thiếu python-docx")
                            doc_obj_x = docx.Document(filepath)
                            
                            # 1. Đọc văn bản thường
                            paras = [p.text for p in doc_obj_x.paragraphs if p.text.strip()]
                            full_text = "\n".join(paras) + "\n\n"
                            
                            # 2. Đọc dữ liệu trong các Bảng biểu (Tables)
                            for table in doc_obj_x.tables:
                                for row in table.rows:
                                    row_data = [cell.text.replace('\n', ' ').strip() for cell in row.cells]
                                    if any(row_data):
                                        full_text += " | ".join(row_data) + "\n"
                                full_text += "\n"
                                
                        elif ext == '.doc':
                            if not win32: raise Exception("Thiếu pywin32")
                            word_app = None
                            doc_obj = None
                            try:
                                word_app = win32.Dispatch("Word.Application")
                                word_app.Visible = False
                                word_app.DisplayAlerts = False
                                doc_obj = word_app.Documents.Open(os.path.abspath(filepath), ReadOnly=True)
                                full_text = doc_obj.Content.Text
                            finally:
                                if doc_obj:
                                    try: doc_obj.Close(False)
                                    except: pass
                                if word_app:
                                    try: word_app.Quit()
                                    except: pass
                    except Exception as e: 
                        log(f"[LOI] Lỗi Word: {e}", "error")
                        continue

                    full_text = re.sub(r'[\x07\x0b\x0c]', '\n', full_text).strip()
                    chunks = [full_text[i:i+8000] for i in range(0, len(full_text), 8000)] if len(full_text) <= 16000 else [full_text[:8000], full_text[-8000:]]
                    if not chunks: chunks = [""] 
                    
                    meta = {"Loại văn bản": "", "Số văn bản": "", "Ngày": "", "Tháng": "", "Năm": "", "Đơn vị soạn văn bản": "", "Người ký": ""}
                    
                    ai_has_error = False
                    for idx, chunk in enumerate(chunks):
                        if not is_running: break
                        for attempt in range(3):
                            if not is_running: break
                            word_prompt = get_prompt_word_metadata(db_conn, doc_types_str) + chunk
                            raw = call_ai_stream(word_prompt, None, idx+1, len(chunks), filename, file_prog_str, force_text_mode=True)
                            ai_data = extract_json(raw, True)
                            if ai_data:
                                for k in ["Loại văn bản", "Số văn bản", "Ngày", "Tháng", "Năm", "Đơn vị soạn văn bản"]:
                                    if not meta.get(k) and ai_data.get(k): meta[k] = str(ai_data.get(k, "")).strip()
                                meta["Người ký"] = merge_signers(meta.get("Người ký", ""), str(ai_data.get("Người ký", "")).strip())
                                ai_has_error = False
                                break
                            else: 
                                log(f"[CANH BAO] Dữ liệu hỏng, thử lại (Lần {attempt+2})...", "warning")
                                ai_has_error = True
                    
                    if is_running:
                        if ai_has_error:
                            log(f"[CANH BAO] Lỗi AI trích xuất Word, đưa vào hàng chờ", "warning")
                            if not is_retry: failed_files.append(data)
                        else:
                            save_document_to_db(db_conn, file_md5, filename, filepath, file_type_str, meta, full_text)
                            log(f"[THANH CONG] Hoàn thành nạp DB Word.", "success")
                            if not is_retry: success_count += 1
                    continue

                is_pdf = ext == '.pdf'
                doc_pdf = fitz.open(filepath) if is_pdf else None
                total_p = len(doc_pdf) if doc_pdf else 1
                meta = {"Loại văn bản": "", "Số văn bản": "", "Ngày": "", "Tháng": "", "Năm": "", "Đơn vị soạn văn bản": "", "Người ký": ""}
                temp_buf = os.path.join(TEMP_DIR, f"buf_{file_md5}.txt")
                with open(temp_buf, "w", encoding="utf-8") as f: f.write("")
                
                has_err = False
                for p_num in range(1, total_p + 1):
                    if not is_running: break
                    tmp_img = os.path.join(TEMP_DIR, f"p{p_num}.jpg")
                    img_ok = False
                    if is_pdf:
                        try:
                            imgs = convert_from_path(filepath, first_page=p_num, last_page=p_num, poppler_path=POPPLER_PATH, dpi=200)
                            if imgs: imgs[0].save(tmp_img, 'JPEG', quality=90); img_ok = True
                        except: pass
                    else:
                        try:
                            img = Image.open(filepath).convert("RGB")
                            img.save(tmp_img, "JPEG", quality=90); img_ok = True
                        except: pass

                    if not img_ok:
                        log(f"[CANH BAO] Lỗi rỗng trang {p_num}", "error"); has_err = True
                        with open(temp_buf, "a", encoding="utf-8") as f: f.write(f"\n--- TRANG {p_num} ---\n[LỖI ẢNH]\n\n")
                        continue

                    with open(tmp_img, "rb") as f_img: b64_img = base64.b64encode(f_img.read()).decode('utf-8')
                    
                    ai_data = None
                    for attempt in range(3):
                        if not is_running: break
                        prompt = get_prompt_page_1(db_conn, doc_types_str) if p_num == 1 else current_prompt_page_n
                        raw = call_ai_stream(prompt, b64_img, p_num, total_p, filename, file_prog_str)
                        ai_data = extract_json(raw, (p_num == 1))
                        if ai_data: break
                        else:
                            if is_running: log(f"[CANH BAO] Dữ liệu JSON bị gãy, đang bóc lại (Lần {attempt+2})...", "warning")

                    if not is_running: break

                    with open(temp_buf, "a", encoding="utf-8") as f:
                        if ai_data:
                            for k in ["Loại văn bản", "Số văn bản", "Ngày", "Tháng", "Năm", "Đơn vị soạn văn bản"]:
                                if not meta.get(k) and ai_data.get(k): meta[k] = str(ai_data.get(k, "")).strip()
                            meta["Người ký"] = merge_signers(meta.get("Người ký", ""), str(ai_data.get("Người ký", "")).strip())
                            tv = ai_data.get("Toàn văn", "")
                            f.write(f"--- TRANG {p_num} ---\n{tv if tv else '[TRỐNG]'}\n\n")
                        else:
                            has_err = True
                            f.write(f"--- TRANG {p_num} ---\n[LỖI AI TỪ CHỐI]\n\n")
                    try: os.remove(tmp_img)
                    except: pass

                if doc_pdf: doc_pdf.close()
                if not is_running: break

                if has_err:
                    log(f"[CANH BAO] File lỗi AI, đưa vào Hàng chờ để thử lại 1 lần cuối.", "warning")
                    if not is_retry: failed_files.append(data)
                else:
                    with open(temp_buf, "r", encoding="utf-8") as f: full_txt = f.read()
                    save_document_to_db(db_conn, file_md5, filename, filepath, file_type_str, meta, full_txt)
                    log(f"[THANH CONG] Đã nạp thành công PDF/Ảnh.", "success")
                    if not is_retry: success_count += 1
                try: os.remove(temp_buf)
                except: pass

        if is_running: run_files(file_data)
        
        if failed_files and is_running:
            log(f"\n{'-'*60}\n[THONG BAO] Phát hiện {len(failed_files)} tập tin bị lỗi AI/Mạng. Tiến hành xử lý lại LẦN CUỐI...\n{'-'*60}", "warning")
            time.sleep(2)
            run_files(failed_files, is_retry=True)
            
        if is_running:
            fail_count = total_files - success_count - skip_count
            log(f"\n{'-'*60}", "system")
            log("[HOAN TAT] HOÀN TẤT TOÀN BỘ QUÁ TRÌNH QUÉT.", "success")
            log(f"[BÁO CÁO] Tổng số File quét: {total_files}", "info")
            log(f"   => Thành công mới: {success_count}", "success")
            log(f"   => Đã có sẵn (Bỏ qua): {skip_count}", "warning")
            log(f"   => Thất bại hoàn toàn: {fail_count}", "error")
            log(f"{'-'*60}\n", "system")

    except Exception as e: 
        if is_running: log(f"\n[LOI FATAL] LỖI HỆ THỐNG NẶNG: {e}", "error")
    finally:
        if db_conn:
            try: db_conn.close()
            except: pass
        if pythoncom:
            try: pythoncom.CoUninitialize()
            except: pass

if __name__ == "__main__":
    log("\n" + "="*60, "system")
    log("[HE THONG] HỆ THỐNG OCR TERMINAL NXBGDHCM", "system")
    log("="*60 + "\n", "system")
    
    if not pptx:
        log("[CANH BAO] Chưa cài thư viện python-pptx. Không thể đọc đuôi .pptx", "warning")
    
    init_global_db()
    load_connections()
    
    doc_types_loaded = False
    if HAS_MYSQL:
        try:
            conn_tmp = get_db_connection()
            if load_document_types(conn_tmp):
                log("[THANH CONG] Đã tải danh mục Loại văn bản từ Máy chủ trung tâm.", "success")
                doc_types_loaded = True
            conn_tmp.close()
        except: pass
        
    target_folder = ""
    if len(sys.argv) > 1:
        target_folder = sys.argv[1].strip('"').strip("'")
        log(f"[KHOI DONG] Kích hoạt tự động với thư mục: {target_folder}", "system")
    else:
        log("\nNhập đường dẫn thư mục cần quét (Hoặc kéo thả thư mục vào đây): ", "info")
        target_folder = input("> ").strip('"').strip("'")

    try:
        if os.path.isdir(target_folder):
            process_batch(target_folder)
        else:
            log(f"[LOI] Đường dẫn không hợp lệ: {target_folder}", "error")
            sys.exit(1)
    except Exception as e:
        log(f"[LOI] Đã có lỗi bất ngờ xảy ra: {e}", "error")